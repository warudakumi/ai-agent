import json
import os
import re
from pathlib import Path

import pandas as pd
import pdfplumber
from app.agent.tools.base import BaseAgentTool
from app.core.error_handler import ErrorSanitizer
from docx import Document
from loguru import logger
from pptx import Presentation


class FileProcessorTool(BaseAgentTool):
    """ファイル処理ツール - 全文読み取り対応"""

    name: str = "file_processor"
    description: str = 'アップロードされたファイルの内容を全文読み取りします。ファイルパスを指定してください。例: {"file_path": "/path/to/file.pdf"}'

    def _run(self, input_str: str) -> str:
        """
        ファイル処理を実行する

        Args:
            input_str: JSON形式の入力。ファイルパスを含む

        Returns:
            ファイルの全文内容
        """
        try:

            def loads_with_windows_path(raw: str) -> dict:
                fixed = re.sub(r'(?<!\\)\\(?![\\"])', r"\\\\", raw)
                return json.loads(fixed)

            inputs = (
                loads_with_windows_path(input_str)
                if isinstance(input_str, str)
                else input_str
            )
            file_path = str(Path(inputs.get("file_path")))

            if not file_path or not os.path.exists(file_path):
                return "エラー: 有効なファイルパスを指定してください"

            # ファイル拡張子の取得
            file_ext = os.path.splitext(file_path)[1].lower()

            # ファイルタイプに応じた処理
            if file_ext == ".pdf":
                return self._process_pdf(file_path)
            elif file_ext == ".docx":
                return self._process_docx(file_path)
            elif file_ext == ".pptx":
                return self._process_pptx(file_path)
            elif file_ext == ".md":
                return self._process_markdown(file_path)
            elif file_ext == ".csv":
                return self._process_csv(file_path)
            elif file_ext == ".txt":
                return self._process_text(file_path)
            elif file_ext in [".xlsx", ".xls"]:
                return self._process_excel(file_path)
            elif file_ext == ".json":
                return self._process_json(file_path)
            else:
                return f"未対応のファイル形式です: {file_ext}"

        except Exception as e:
            logger.error(f"ファイル処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_pdf(self, file_path: str) -> str:
        """PDFファイルを処理"""
        try:
            content = "=== PDFファイル内容 ===\n\n"

            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        content += f"--- ページ {page_num} ---\n"
                        content += text.strip() + "\n\n"

                    # テーブルがある場合は抽出
                    tables = page.extract_tables()
                    for table_num, table in enumerate(tables, 1):
                        if table:
                            content += f"[ページ {page_num} - テーブル {table_num}]\n"
                            for row in table:
                                if row:
                                    content += (
                                        " | ".join([cell or "" for cell in row]) + "\n"
                                    )
                            content += "\n"

            return (
                content
                if content.strip() != "=== PDFファイル内容 ==="
                else "PDFからテキストを抽出できませんでした"
            )

        except Exception as e:
            logger.error(f"PDF処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_docx(self, file_path: str) -> str:
        """Wordファイルを処理"""
        try:
            doc = Document(file_path)
            content = "=== Wordファイル内容 ===\n\n"

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # 見出しレベルの判定
                    if paragraph.style.name.startswith("Heading"):
                        level = paragraph.style.name.replace("Heading ", "")
                        if level.isdigit():
                            content += "#" * int(level) + " " + paragraph.text + "\n\n"
                        else:
                            content += "# " + paragraph.text + "\n\n"
                    else:
                        content += paragraph.text + "\n\n"

            # テーブルの処理
            for table_num, table in enumerate(doc.tables, 1):
                content += f"[テーブル {table_num}]\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        content += row_text + "\n"
                content += "\n"

            return (
                content
                if content.strip() != "=== Wordファイル内容 ==="
                else "Wordファイルからテキストを抽出できませんでした"
            )

        except Exception as e:
            logger.error(f"Word処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_pptx(self, file_path: str) -> str:
        """PowerPointファイルを処理"""
        try:
            prs = Presentation(file_path)
            content = "=== PowerPointファイル内容 ===\n\n"

            for slide_num, slide in enumerate(prs.slides, 1):
                content += f"--- スライド {slide_num} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # タイトルかどうかの判定
                        if shape == slide.shapes.title:
                            content += f"# {shape.text}\n\n"
                        else:
                            content += f"{shape.text}\n\n"

                    # テーブルの処理
                    if shape.has_table:
                        content += "[テーブル]\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                [cell.text.strip() for cell in row.cells]
                            )
                            if row_text.strip():
                                content += row_text + "\n"
                        content += "\n"

                content += "\n"

            return (
                content
                if content.strip() != "=== PowerPointファイル内容 ==="
                else "PowerPointファイルからテキストを抽出できませんでした"
            )

        except Exception as e:
            logger.error(f"PowerPoint処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_markdown(self, file_path: str) -> str:
        """Markdownファイルを処理"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            # Markdownの構造をそのまま保持
            result = "=== Markdownファイル内容 ===\n\n"
            result += content

            return result

        except Exception as e:
            logger.error(f"Markdown処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_csv(self, file_path: str) -> str:
        """CSVファイルを処理"""
        try:
            df = pd.read_csv(file_path)

            content = "=== CSVファイル内容 ===\n\n"
            content += f"レコード数: {len(df)}\n"
            content += f"カラム数: {len(df.columns)}\n\n"

            # 全データを表形式で出力
            content += df.to_string(index=False)

            return content

        except Exception as e:
            logger.error(f"CSV処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_text(self, file_path: str) -> str:
        """テキストファイルを処理"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            result = "=== テキストファイル内容 ===\n\n"
            result += content

            return result

        except Exception as e:
            logger.error(f"テキスト処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_excel(self, file_path: str) -> str:
        """Excelファイルを処理"""
        try:
            excel = pd.ExcelFile(file_path)
            content = "=== Excelファイル内容 ===\n\n"

            for sheet_name in excel.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content += f"--- シート: {sheet_name} ---\n"
                content += f"レコード数: {len(df)}\n"
                content += f"カラム数: {len(df.columns)}\n\n"

                # 全データを表形式で出力
                content += df.to_string(index=False)
                content += "\n\n"

            return content

        except Exception as e:
            logger.error(f"Excel処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message

    def _process_json(self, file_path: str) -> str:
        """JSONファイルを処理"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            content = "=== JSONファイル内容 ===\n\n"
            content += json.dumps(data, indent=2, ensure_ascii=False)

            return content

        except Exception as e:
            logger.error(f"JSON処理エラー: {str(e)}")
            safe_message = ErrorSanitizer.sanitize_error_message(
                str(e), "file_processing"
            )
            return safe_message
